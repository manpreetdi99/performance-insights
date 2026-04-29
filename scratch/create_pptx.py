from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Color palette ──────────────────────────────────────────────────────────────
BG        = RGBColor(0x0F, 0x17, 0x23)   # dark navy
CARD      = RGBColor(0x1A, 0x24, 0x33)   # card background
RED       = RGBColor(0xEF, 0x44, 0x44)   # brand red
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
MUTED     = RGBColor(0x94, 0xA3, 0xB8)
ACCENT    = RGBColor(0x38, 0xBD, 0xF8)   # sky blue
GREEN     = RGBColor(0x22, 0xC5, 0x5E)
YELLOW    = RGBColor(0xFA, 0xCC, 0x15)
VIOLET    = RGBColor(0xA7, 0x8B, 0xFA)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank layout

# ── Helpers ────────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def slide_bg(slide):
    add_rect(slide, 0, 0, 13.33, 7.5, BG)

def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.1, CARD)
    add_rect(slide, 0, 1.1, 13.33, 0.04, RED)
    add_text(slide, title, 0.4, 0.15, 10, 0.6, size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.7, 10, 0.4, size=12, color=MUTED)

def bullet_card(slide, l, t, w, h, title, bullets, title_color=ACCENT):
    add_rect(slide, l, t, w, h, CARD)
    # top accent line
    add_rect(slide, l, t, w, 0.035, title_color)
    add_text(slide, title, l+0.15, t+0.08, w-0.3, 0.35, size=13, bold=True, color=title_color)
    bullet_top = t + 0.45
    for b in bullets:
        add_text(slide, f"▸  {b}", l+0.15, bullet_top, w-0.3, 0.3, size=10.5, color=WHITE)
        bullet_top += 0.29

def stat_box(slide, l, t, number, label, color=ACCENT):
    add_rect(slide, l, t, 2.8, 1.4, CARD)
    add_rect(slide, l, t, 2.8, 0.04, color)
    add_text(slide, number, l, t+0.15, 2.8, 0.65, size=32, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, label,  l, t+0.75, 2.8, 0.4,  size=11, color=MUTED, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – Title
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)

# diagonal accent strip
add_rect(sl, 8.2, 0, 5.13, 7.5, CARD)
add_rect(sl, 8.18, 0, 0.06, 7.5, RED)

add_text(sl, "FASMETRICS", 0.6, 1.6, 7.5, 1.2, size=52, bold=True, color=WHITE)
add_text(sl, "Analytics", 0.6, 2.7, 7.5, 0.9, size=42, bold=True, color=RED)
add_text(sl, "Network Quality Benchmarking Platform",
         0.6, 3.6, 7.4, 0.6, size=16, color=MUTED)
add_text(sl, "End-to-end mobile call analysis · LTE · GSM · SRVCC",
         0.6, 4.2, 7.4, 0.5, size=13, color=ACCENT)

# right panel labels
for i, (label, val) in enumerate([
    ("Platform", "Web + REST API"),
    ("Database", "SQL Server"),
    ("Frontend", "React + TypeScript"),
    ("Backend", "Python FastAPI"),
]):
    add_text(sl, label, 8.5, 1.8+i*0.9, 2.0, 0.35, size=10, color=MUTED)
    add_text(sl, val,   8.5, 2.1+i*0.9, 4.4, 0.4,  size=13, bold=True, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – Problem & Solution
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "Problem & Solution", "Why FASMETRICS exists")

# Problem
add_rect(sl, 0.4, 1.3, 5.9, 5.6, CARD)
add_rect(sl, 0.4, 1.3, 5.9, 0.04, RED)
add_text(sl, "❌  The Problem", 0.6, 1.35, 5.5, 0.45, size=14, bold=True, color=RED)
problems = [
    "Manual log review is slow & error-prone",
    "No unified view across multiple test sessions",
    "Difficult to correlate GSM/LTE signal drops with call failures",
    "No quick way to mark & filter invalid / 'fake' test calls",
    "Radio measurements, KPIs and trace logs exist in silos",
    "Engineers spend hours in raw SQL to answer simple questions",
]
for i, p in enumerate(problems):
    add_text(sl, f"•  {p}", 0.6, 1.95+i*0.72, 5.5, 0.55, size=11, color=WHITE)

# Solution
add_rect(sl, 6.7, 1.3, 6.2, 5.6, CARD)
add_rect(sl, 6.7, 1.3, 6.2, 0.04, GREEN)
add_text(sl, "✅  The Solution", 6.9, 1.35, 5.8, 0.45, size=14, bold=True, color=GREEN)
solutions = [
    "Centralised dashboard with live SQL data",
    "Unified call list across collections & databases",
    "Time-synced radio charts with hover highlighting",
    "Automatic invalid-session detection via comment",
    "Side-by-side A/B signal measurement comparison",
    "Built-in query editor for ad-hoc SQL benchmarks",
]
for i, s in enumerate(solutions):
    add_text(sl, f"•  {s}", 6.9, 1.95+i*0.72, 5.8, 0.55, size=11, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – Key Features Overview
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "Key Features", "Six pillars of the platform")

cards = [
    (0.35, 1.25, "📋  Call Management",
     ["Multi-database & collection selector",
      "Global filter panel (accessible from any tab)",
      "Status / Session-valid / Location filters",
      "Color-coded rows (drop · fail · invalid · release)"], ACCENT),
    (4.7,  1.25, "📡  Radio Analysis",
     ["LTE: RSRP / RSRQ time-series charts",
      "GSM: RxLev / RxQual with thresholds",
      "A-side & B-side toggle",
      "SRVCC: LTE → GSM handover view"], GREEN),
    (9.05, 1.25, "🗺  Interactive Map",
     ["Calls plotted on GPS coordinates",
      "Click a pin to jump to Call Detail",
      "Filter by location & device"], YELLOW),
    (0.35, 4.15, "🏷  Smart Comments",
     ["Quick-select dropdown (LC/LQ/FAKE…)",
      "FAKE prefix → auto set Valid = 0",
      "Change to non-FAKE → auto restore Valid = 1",
      "Persisted to AnalysisCommentSessionsBridge"], RED),
    (4.7,  4.15, "📊  KPI & TraceLog",
     ["ResultsKPI table with chrono ordering",
      "FactSystemTraceLog for A+B sessions",
      "Hover a row → highlight on signal chart",
      "MOS tooltip with individual values"], VIOLET),
    (9.05, 4.15, "🛠  Query Editor",
     ["Run ad-hoc SQL on any connected DB",
      "Multi-query benchmark mode",
      "Execution time & row count per query",
      "Collection-aware query templates"], ACCENT),
]

for (l, t, title, bullets, color) in cards:
    bullet_card(sl, l, t, 4.0, 2.7, title, bullets, title_color=color)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – Architecture
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "Architecture", "Full-stack overview")

layers = [
    (ACCENT,  "Frontend",  "React 18 · TypeScript · Vite · shadcn/ui · Recharts · Framer Motion · Leaflet"),
    (GREEN,   "API Layer", "Python FastAPI · pyodbc · REST endpoints · CORS middleware"),
    (YELLOW,  "Database",  "Microsoft SQL Server · Sessions · CallAnalysis · LTE/GSM Measurements · TraceLog · KPI"),
]
for i, (color, title, desc) in enumerate(layers):
    y = 1.5 + i*1.65
    add_rect(sl, 1.0, y, 11.33, 1.3, CARD)
    add_rect(sl, 1.0, y, 0.07, 1.3, color)
    add_text(sl, title, 1.3, y+0.1,  3.0, 0.45, size=15, bold=True, color=color)
    add_text(sl, desc,  1.3, y+0.55, 10.7, 0.65, size=12, color=WHITE)
    # arrows between layers
    if i < 2:
        add_text(sl, "▼", 6.5, y+1.3, 0.5, 0.35, size=16, color=MUTED, align=PP_ALIGN.CENTER)

# Data flow note
add_rect(sl, 1.0, 6.55, 11.33, 0.65, CARD)
add_text(sl, "Data flow:  Browser → Vite Dev Server → FastAPI (localhost:8000) → SQL Server",
         1.2, 6.6, 11.0, 0.5, size=11, color=MUTED)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – Call Detail Deep Dive
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "Call Detail View", "Everything about one call in one screen")

features = [
    (ACCENT, "Signal Chart",
     ["RSRP / RSRQ (LTE) or RxLev / RxQual (GSM)",
      "Show/hide individual series via checkboxes",
      "Warning & critical reference lines",
      "Hover a table row → vertical marker on chart"]),
    (GREEN, "A/B Side Comparison",
     ["Side-by-side call status table (A vs B)",
      "Switch between A-side and B-side measurements",
      "SRVCC: switch between LTE and GSM networks",
      "B-side LTE summary statistics"]),
    (YELLOW, "KPI Results",
     ["All KPI records for the session",
      "StartTime, KPIId, ErrorCode, Values 3-5",
      "Hover row → chart reference line sync",
      "Scrollable, sticky header"]),
    (VIOLET, "TraceLog",
     ["FactSystemTraceLog for A+B sessions",
      "FullDate · Side · SessionId · Info columns",
      "Hover sync with signal chart",
      "Colour highlight on active row"]),
]

for i, (color, title, bullets) in enumerate(features):
    col = i % 2
    row = i // 2
    bullet_card(sl, 0.35 + col*6.5, 1.25 + row*2.9, 6.15, 2.7, title, bullets, title_color=color)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – Fake / Invalid Call Workflow
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "Smart Call Validation", "Automatic Valid/Invalid management via comments")

steps = [
    (RED,   "1. Select Comment",  "Engineer picks a FAKE option\n(FAKE UE STUCK · FAKE NO SYNC · FAKE EOF)\nor types any comment starting with 'fake'"),
    (YELLOW,"2. Save",           "Frontend calls POST /api/calls/comment\nwith session_id + comment text"),
    (RED,   "3. Backend Check",  "Python detects comment.lower().startswith('fake')\n→ UPDATE Sessions SET Valid = 0"),
    (GREEN, "4. Re-validation",  "If comment changed to non-fake (e.g. LC LTE)\n→ UPDATE Sessions SET Valid = 1 automatically"),
]

for i, (color, title, desc) in enumerate(steps):
    add_rect(sl, 0.4 + i*3.2, 1.5, 3.0, 3.5, CARD)
    add_rect(sl, 0.4 + i*3.2, 1.5, 3.0, 0.05, color)
    add_text(sl, title, 0.55+i*3.2, 1.6,  2.7, 0.5, size=13, bold=True, color=color)
    add_text(sl, desc,  0.55+i*3.2, 2.2,  2.7, 2.5, size=11, color=WHITE, wrap=True)
    if i < 3:
        add_text(sl, "→", 3.3+i*3.2, 2.9, 0.5, 0.5, size=22, color=MUTED, align=PP_ALIGN.CENTER)

# Result banner
add_rect(sl, 0.4, 5.3, 12.5, 1.0, CARD)
add_rect(sl, 0.4, 5.3, 12.5, 0.05, ACCENT)
add_text(sl,
    "Result: Call list auto-highlights invalid rows in red · Filter panel shows Valid/Invalid toggle · "
    "No manual DB edits required",
    0.6, 5.4, 12.2, 0.7, size=12, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – Advantages & Benefits
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
header_bar(sl, "Key Advantages", "Why this tool is better than the alternatives")

advantages = [
    (ACCENT, "⚡ Speed",
     ["Instant call loading with async API",
      "Local SQL Server = zero network latency",
      "Optimised queries with ordered indexes"]),
    (GREEN, "🎯 Accuracy",
     ["Direct DB reads — no data export step",
      "Time-synced chart ↔ table interaction",
      "Automatic session validation logic"]),
    (YELLOW, "🔒 Control",
     ["Full control over database & queries",
      "No third-party cloud dependency",
      "Query editor for custom SQL analysis"]),
    (RED, "👁 Visibility",
     ["Colour-coded rows for instant triage",
      "End-of-File markers between log files",
      "Global filter badge shows active state"]),
    (VIOLET, "🔄 Flexibility",
     ["Multi-database support in one tool",
      "GSM / LTE / SRVCC call modes",
      "A-side & B-side independent views"]),
    (ACCENT, "🚀 Workflow",
     ["Single-click from call list to detail",
      "Filter panel accessible from any tab",
      "Persistent settings via localStorage"]),
]

for i, (color, title, bullets) in enumerate(advantages):
    col = i % 3
    row = i // 3
    bullet_card(sl, 0.35 + col*4.3, 1.25 + row*2.85, 4.1, 2.65, title, bullets, title_color=color)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – Summary / Thank You
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
slide_bg(sl)
add_rect(sl, 0, 0, 13.33, 7.5, BG)
add_rect(sl, 0, 0, 13.33, 0.06, RED)
add_rect(sl, 0, 7.44, 13.33, 0.06, RED)

add_text(sl, "FASMETRICS Analytics", 1.5, 1.5, 10, 1.0, size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "One platform · All calls · Full insight",
         1.5, 2.55, 10, 0.55, size=18, color=RED, align=PP_ALIGN.CENTER)

summary_points = [
    "✔  Live SQL data — no exports, no delays",
    "✔  Unified view: LTE · GSM · SRVCC · A/B sides",
    "✔  Interactive charts with hover sync",
    "✔  Smart comment-based call validation",
    "✔  Global filter panel accessible from every tab",
    "✔  Built-in query editor for ad-hoc analysis",
]
for i, pt in enumerate(summary_points):
    add_text(sl, pt, 2.5, 3.4+i*0.52, 8.5, 0.45, size=13, color=WHITE, align=PP_ALIGN.CENTER)

add_text(sl, "Stack: React · TypeScript · Python FastAPI · SQL Server",
         1.5, 6.8, 10, 0.45, size=11, color=MUTED, align=PP_ALIGN.CENTER)

# ── Save ────────────────────────────────────────────────────────────────────
out = r"c:\Users\Mechanical Engineer\Documents\performance-insights\FASMETRICS_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
